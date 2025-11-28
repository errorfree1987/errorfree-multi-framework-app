"""
Error-Free® Multi-framework AI Document Analyzer
- Internal high-level accounts (admin/pro/test): GPT-5.1, unlimited
- Guest accounts: GPT-4.1-mini, 文件 3 份 / 每框架 1 次分析 / 3 次追問 / 1 次下載
- 支援 PDF / DOCX / TXT / 圖片(JPG/PNG，含簡單 OCR)
- 多框架並存：每個框架的分析結果與 Q&A 保留，不互相覆蓋
"""

import os
import json
import base64
import datetime
from pathlib import Path
from typing import Dict, List, Tuple
from io import BytesIO

import streamlit as st
import pdfplumber
from docx import Document
from openai import OpenAI
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.pagesizes import letter

# PPTX（簡單匯出）
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# =========================
# 檔案位置
# =========================

BASE_DIR = Path(".")
GUEST_FILE = BASE_DIR / "guest_accounts.json"
STATE_FILE = BASE_DIR / "user_state.json"
DOC_TRACK_FILE = BASE_DIR / "user_docs.json"
USAGE_FILE = BASE_DIR / "usage_stats.json"

# =========================
# 帳號設定
# =========================

ACCOUNTS = {
    "admin@errorfree.com": {"password": "1111", "role": "admin"},
    "dr.chiu@errorfree.com": {"password": "2222", "role": "pro"},
    "test@errorfree.com": {"password": "3333", "role": "pro"},
}

# =========================
# OpenAI Client
# =========================

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None


# =========================
# JSON helpers
# =========================
def _load_json(path: Path, default):
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def _save_json(path: Path, data):
    try:
        path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        pass


def load_guest_accounts() -> Dict[str, Dict]:
    return _load_json(GUEST_FILE, {})


def save_guest_accounts(data: Dict[str, Dict]):
    _save_json(GUEST_FILE, data)


def load_doc_tracking() -> Dict[str, List[str]]:
    return _load_json(DOC_TRACK_FILE, {})


def save_doc_tracking(data: Dict[str, List[str]]):
    _save_json(DOC_TRACK_FILE, data)


def load_usage_stats() -> Dict[str, Dict]:
    return _load_json(USAGE_FILE, {})


def save_usage_stats(data: Dict[str, Dict]):
    _save_json(USAGE_FILE, data)


# =========================
# Framework 定義
# =========================

FRAMEWORKS: Dict[str, Dict] = {
    "omission": {
        "name_zh": "Error-Free® 遺漏錯誤檢查框架",
        "name_en": "Error-Free® Omission Error Check Framework",
        "wrapper_zh": (
            "你是一位 Error-Free® 遺漏錯誤檢查專家。"
            "請分析文件中可能遺漏的重要內容、條件、假設、角色、步驟、風險或例外，"
            "並說明遺漏的影響與具體補強建議，最後整理成條列與一個簡單的 Markdown 表格。"
            "請盡量使用一般的 '-' 或數字條列，不要使用特殊符號（例如 ■ 或 ●）。"
        ),
        "wrapper_en": (
            "You are an Error-Free® omission error expert. "
            "Review the document, find important missing information or conditions, "
            "explain the impact, and give concrete suggestions. "
            "Use plain '-' or numbered lists (no special bullets like ■ or ●). "
            "Finish with a simple Markdown table."
        ),
    },
    "technical": {
        "name_zh": "Error-Free® 技術風險檢查框架",
        "name_en": "Error-Free® Technical Risk Check Framework",
        "wrapper_zh": (
            "你是一位 Error-Free® 技術風險檢查專家。"
            "請從技術假設、邊界條件、相容性、安全性、可靠度與單點失敗等面向分析文件，"
            "列出技術風險、風險等級與實務改善建議，並以 Markdown 表格整理重點。"
            "請盡量使用 '-' 或數字條列，不要使用特殊符號（例如 ■ 或 ●）。"
        ),
        "wrapper_en": (
            "You are an Error-Free® technical risk review expert. "
            "Analyze the document for technical assumptions, edge cases, compatibility, "
            "safety and single points of failure. List risks, risk level and mitigation. "
            "Use plain '-' or numbered lists only (no special bullets). "
            "Provide a summary Markdown table."
        ),
    },
}


# =========================
# State 儲存
# =========================
def save_state_to_disk():
    data = {
        "user_email": st.session_state.get("user_email"),
        "user_role": st.session_state.get("user_role"),
        "is_authenticated": st.session_state.get("is_authenticated", False),
        "lang": st.session_state.get("lang", "zh"),
        "last_doc_text": st.session_state.get("last_doc_text", ""),
        "framework_states": st.session_state.get("framework_states", {}),
        "selected_framework_key": st.session_state.get("selected_framework_key"),
        "current_doc_id": st.session_state.get("current_doc_id"),
    }
    _save_json(STATE_FILE, data)


def restore_state_from_disk():
    data = _load_json(STATE_FILE, {})
    for k, v in data.items():
        if k not in st.session_state:
            st.session_state[k] = v


# =========================
# Usage 統計（給 admin 看）
# =========================
def record_usage(user_email: str, framework_key: str, kind: str):
    if not user_email:
        return
    data = load_usage_stats()
    u = data.get(user_email, {})
    fw_map = u.get("frameworks", {})
    fw_entry = fw_map.get(
        framework_key, {"analysis_runs": 0, "followups": 0, "downloads": 0}
    )
    if kind == "analysis":
        fw_entry["analysis_runs"] = fw_entry.get("analysis_runs", 0) + 1
    elif kind == "followup":
        fw_entry["followups"] = fw_entry.get("followups", 0) + 1
    elif kind == "download":
        fw_entry["downloads"] = fw_entry.get("downloads", 0) + 1
    fw_map[framework_key] = fw_entry
    u["frameworks"] = fw_map
    u["last_used"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    data[user_email] = u
    save_usage_stats(data)


# =========================
# Model routing
# =========================
def resolve_model_for_user(role: str) -> str:
    if role in ["admin", "pro"]:
        return "gpt-5.1"
    return "gpt-4.1-mini"  # guest / others


# =========================
# 檔案讀取 & OCR
# =========================
def clean_binary_image(uploaded_file) -> Tuple[bytes, str]:
    data = uploaded_file.read()
    name = uploaded_file.name.lower()
    if name.endswith(".png"):
        mime = "image/png"
    else:
        mime = "image/jpeg"
    return data, mime


def ocr_image_to_text(uploaded_file) -> str:
    """使用 OpenAI 做簡單 OCR，把圖片裡文字抽出來。"""
    if client is None:
        return f"[圖片：{uploaded_file.name}]（因未設定 OPENAI_API_KEY，無法 OCR）"

    data, mime = clean_binary_image(uploaded_file)
    b64 = base64.b64encode(data).decode("utf-8")
    image_url = f"data:{mime};base64,{b64}"

    try:
        resp = client.responses.create(
            model="gpt-4.1-mini",
            input=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "input_text",
                            "text": (
                                "Please read all visible text from this image and return "
                                "a clean plain-text transcription. Do not explain."
                            ),
                        },
                        {
                            "type": "input_image",
                            "image_url": {"url": image_url},
                        },
                    ],
                }
            ],
            max_output_tokens=1200,
        )
        text = resp.output_text
        return f"[圖片 OCR 文字摘錄，自 {uploaded_file.name}]\n\n" + text
    except Exception as e:
        return f"[圖片：{uploaded_file.name}]（OCR 失敗：{e}）"


def read_file_to_text(uploaded_file) -> str:
    """支援 PDF / DOCX / TXT / 圖片（JPG/PNG + OCR）"""
    if uploaded_file is None:
        return ""
    name = uploaded_file.name.lower()

    try:
        if name.endswith(".pdf"):
            text_pages: List[str] = []
            with pdfplumber.open(uploaded_file) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    text_pages.append(t)
            return "\n".join(text_pages)
        elif name.endswith(".docx"):
            doc = Document(uploaded_file)
            return "\n".join(p.text for p in doc.paragraphs)
        elif name.endswith(".txt"):
            return uploaded_file.read().decode("utf-8", errors="ignore")
        elif name.endswith((".jpg", ".jpeg", ".png")):
            return ocr_image_to_text(uploaded_file)
        else:
            return f"[不支援的檔案類型：{uploaded_file.name}]"
    except Exception as e:
        return f"[讀取檔案時發生錯誤: {e}]"


# =========================
# LLM 呼叫
# =========================
def run_llm_analysis(
    framework_key: str, language: str, document_text: str, model_name: str
) -> str:
    fw = FRAMEWORKS[framework_key]
    system_prompt = fw["wrapper_zh"] if language == "zh" else fw["wrapper_en"]
    prefix = "以下是要分析的文件內容：\n\n" if language == "zh" else "Here is the document to analyze:\n\n"
    user_prompt = prefix + document_text

    if client is None:
        return "[Error] OPENAI_API_KEY 尚未設定，無法連線至 OpenAI。"

    try:
        resp = client.responses.create(
            model=model_name,
            input=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            max_output_tokens=2500,
        )
        return resp.output_text
    except Exception as e:
        return f"[呼叫 OpenAI API 時發生錯誤: {e}]"


def run_followup_qa(
    framework_key: str,
    language: str,
    document_text: str,
    analysis_output: str,
    user_question: str,
    model_name: str,
    extra_text: str = "",
) -> str:
    fw = FRAMEWORKS[framework_key]
    if language == "zh":
        system_prompt = (
            "You are an Error-Free consultant familiar with framework: "
            + fw["name_zh"]
            + "。你已經對文件完成一次完整分析，現在只針對追問補充說明，避免重複整份報告。"
            "可引用原文件、先前分析以及補充附件內容。"
        )
    else:
        system_prompt = (
            "You are an Error-Free consultant for framework: "
            + fw["name_en"]
            + ". You already produced a full analysis. "
            "Answer follow-up questions based on the original document, previous analysis and extra attachments, "
            "and avoid recreating the full report."
        )

    doc_excerpt = document_text[:8000]
    analysis_excerpt = analysis_output[:8000]
    extra_excerpt = extra_text[:4000] if extra_text else ""

    blocks = [
        "Original document excerpt:\n" + doc_excerpt,
        "Previous analysis excerpt:\n" + analysis_excerpt,
        "User question:\n" + user_question,
    ]
    if extra_excerpt:
        blocks.append("Extra reference:\n" + extra_excerpt)

    user_content = "\n\n".join(blocks)

    if client is None:
        return "[Error] OPENAI_API_KEY 尚未設定。"

    try:
        resp = client.responses.create(
            model=model_name,
            input=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content},
            ],
            max_output_tokens=2000,
        )
        return resp.output_text
    except Exception as e:
        return f"[呼叫 OpenAI API 時發生錯誤: {e}]"


# =========================
# 報告清理 & 匯出
# =========================
def clean_report_text(text: str) -> str:
    """把容易變成黑方塊的符號換掉，避免 PDF 裡亂碼。"""
    replacements = {
        "■": "-",
        "●": "-",
        "▪": "-",
        "◼": "-",
        "•": "-",
        "–": "-",
        "—": "-",
        "\u2022": "-",
        "\u25cf": "-",
        "\u25a0": "-",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def build_full_report(lang: str, framework_key: str, state: Dict) -> str:
    analysis_output = state.get("analysis_output", "")
    followups = state.get("followup_history", [])
    fw = FRAMEWORKS[framework_key]
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    email = st.session_state.get("user_email", "unknown")

    if lang == "zh":
        parts: List[str] = [
            "Error-Free® 多框架 AI 文件分析報告（分析 + Q&A）",
            f"產生時間：{now}",
            f"使用者帳號：{email}",
            f"使用框架：{fw['name_zh']}",
            "",
            "==============================",
            "一、分析結果",
            "==============================",
            analysis_output,
        ]
        if followups:
            parts += [
                "",
                "==============================",
                "二、後續問答（Q&A）",
                "==============================",
            ]
            for i, (q, a) in enumerate(followups, start=1):
                parts.append(f"[Q{i}] {q}")
                parts.append(f"[A{i}] {a}")
                parts.append("")
    else:
        parts = [
            "Error-Free® Multi-framework AI Report (Analysis + Q&A)",
            f"Generated: {now}",
            f"User: {email}",
            f"Framework: {fw['name_en']}",
            "",
            "==============================",
            "1. Analysis",
            "==============================",
            analysis_output,
        ]
        if followups:
            parts += [
                "",
                "==============================",
                "2. Follow-up Q&A",
                "==============================",
            ]
            for i, (q, a) in enumerate(followups, start=1):
                parts.append(f"[Q{i}] {q}")
                parts.append(f"[A{i}] {a}")
                parts.append("")

    raw = "\n".join(parts)
    return clean_report_text(raw)


def build_docx_bytes(text: str) -> bytes:
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def build_pdf_bytes(text: str) -> bytes:
    buf = BytesIO()
    c = pdf_canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    y = height - 40

    for line in text.split("\n"):
        safe_line = line.replace("\t", "    ")
        c.drawString(40, y, safe_line[:1000])
        y -= 14
        if y < 40:
            c.showPage()
            y = height - 40
    c.save()
    buf.seek(0)
    return buf.getvalue()


def build_pptx_bytes(text: str) -> bytes:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx 未安裝，無法匯出 PPTX。")

    prs = Presentation()
    # 第一頁：標題
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Error-Free® Report"
    slide.placeholders[1].text = "Analysis + Q&A"

    # 之後的文字切成多頁，每頁最多 12 行
    lines = [l for l in text.split("\n") if l.strip()]
    chunk_size = 12
    for i in range(0, len(lines), chunk_size):
        chunk = lines[i : i + chunk_size]
        layout = prs.slide_layouts[1]  # Title + Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = f"Section {i // chunk_size + 1}"
        body = s.placeholders[1]
        tf = body.text_frame
        tf.text = chunk[0][:200]
        for line in chunk[1:]:
            p = tf.add_paragraph()
            p.text = line[:200]
            p.level = 0

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# =========================
# Admin Dashboard（只給 admin/pro 看）
# =========================
def admin_dashboard():
    lang = st.session_state.get("lang", "zh")
    st.title("Admin Dashboard — Error-Free®")
    st.markdown("---")

    # Guest accounts
    st.subheader("📌 Guest 帳號列表" if lang == "zh" else "📌 Guest accounts")
    guests = load_guest_accounts()
    if not guests:
        st.info("目前沒有 Guest 帳號。" if lang == "zh" else "No guest accounts yet.")
    else:
        for email, acc in guests.items():
            st.markdown(f"**{email}** — password: `{acc.get('password')}`")
        st.markdown("---")

    # Guest documents
    st.subheader("📁 Guest 文件使用狀況" if lang == "zh" else "📁 Guest document usage")
    doc_tracking = load_doc_tracking()
    if not doc_tracking:
        st.info(
            "尚無 Guest 上傳記錄。" if lang == "zh" else "No guest uploads recorded yet."
        )
    else:
        for email, docs in doc_tracking.items():
            st.markdown(f"**{email}** — 上傳文件數：{len(docs)} / 3")
            for d in docs:
                st.markdown(f"- {d}")
            st.markdown("---")

    # Current session framework state
    st.subheader(
        "🧩 模組分析與追問狀況 (目前 Session)"
        if lang == "zh"
        else "🧩 Framework state (current session)"
    )
    fs = st.session_state.get("framework_states", {})
    if not fs:
        st.info(
            "尚無 Framework 分析記錄" if lang == "zh" else "No framework analysis yet."
        )
    else:
        for fw_key, state in fs.items():
            fw_name = (
                FRAMEWORKS[fw_key]["name_zh"]
                if lang == "zh"
                else FRAMEWORKS[fw_key]["name_en"]
            )
            st.markdown(f"### ▶ {fw_name}")
            st.write(f"分析完成：{state.get('analysis_done')}")
            st.write(f"追問次數：{len(state.get('followup_history', []))}")
            st.write(f"已下載報告：{state.get('download_used')}")
            st.markdown("---")

    if st.button("返回分析頁面" if lang == "zh" else "Back to analysis"):
        st.session_state.show_admin = False
        save_state_to_disk()
        st.experimental_rerun()


# =========================
# Main App
# =========================
def language_selector():
    current = st.session_state.get("lang", "zh")
    index = 0 if current == "en" else 1
    choice = st.radio("Language / 語言", ("English", "中文"), index=index)
    st.session_state.lang = "en" if choice == "English" else "zh"


def main():
    st.set_page_config(
        page_title="Error-Free® Multi-framework Analyzer", layout="wide"
    )
    restore_state_from_disk()

    for k, v in [
        ("user_email", None),
        ("user_role", None),
        ("is_authenticated", False),
        ("lang", "zh"),
        ("last_doc_text", ""),
        ("framework_states", {}),
        ("selected_framework_key", list(FRAMEWORKS.keys())[0]),
        ("current_doc_id", None),
        ("show_admin", False),
    ]:
        if k not in st.session_state:
            st.session_state[k] = v

    # Sidebar
    with st.sidebar:
        lang = st.session_state.lang
        language_selector()

        if (
            st.session_state.is_authenticated
            and st.session_state.user_role in ["admin", "pro"]
        ):
            if st.button("管理後台 Admin Dashboard"):
                st.session_state.show_admin = True
                save_state_to_disk()
                st.experimental_rerun()

        st.markdown("---")
        if st.session_state.is_authenticated:
            st.subheader("帳號資訊" if lang == "zh" else "Account")
            st.write(f"Email：{st.session_state.user_email}")
            role = st.session_state.user_role
            st.write(
                "角色：內部帳號"
                if lang == "zh" and role in ["admin", "pro"]
                else ("角色：Guest" if lang == "zh" else f"Role: {role}")
            )
            if st.button("登出" if lang == "zh" else "Logout"):
                st.session_state.user_email = None
                st.session_state.user_role = None
                st.session_state.is_authenticated = False
                st.session_state.framework_states = {}
                st.session_state.last_doc_text = ""
                save_state_to_disk()
                st.experimental_rerun()
        else:
            st.subheader("尚未登入" if lang == "zh" else "Not logged in")

    # Admin dashboard
    if st.session_state.show_admin and st.session_state.is_authenticated:
        if st.session_state.user_role in ["admin", "pro"]:
            admin_dashboard()
            return
        else:
            st.session_state.show_admin = False

    # Login Page
    if not st.session_state.is_authenticated:
        lang = st.session_state.lang
        title = (
            "Error-Free® 多框架文件分析"
            if lang == "zh"
            else "Error-Free® Multi-framework Document Analyzer"
        )
        st.title(title)
        st.markdown("---")

        col_emp, col_guest = st.columns(2)

        # Internal login
        with col_emp:
            st.markdown(
                "### 內部員工 / 會員登入"
                if lang == "zh"
                else "### Internal Employee / Member Login"
            )
            emp_email = st.text_input("Email", key="emp_email")
            emp_pw = st.text_input(
                "密碼" if lang == "zh" else "Password",
                type="password",
                key="emp_pw",
            )
            if st.button("登入" if lang == "zh" else "Login", key="emp_login_btn"):
                account = ACCOUNTS.get(emp_email)
                if account and account["password"] == emp_pw:
                    st.session_state.user_email = emp_email
                    st.session_state.user_role = account["role"]
                    st.session_state.is_authenticated = True
                    save_state_to_disk()
                    st.experimental_rerun()
                else:
                    st.error(
                        "帳號或密碼錯誤"
                        if lang == "zh"
                        else "Invalid email or password"
                    )

        # Guest login/signup
        with col_guest:
            st.markdown("### Guest 試用" if lang == "zh" else "### Guest Trial")

            st.markdown("**Guest 登入**" if lang == "zh" else "**Guest Login**")
            g_email = st.text_input("Guest Email", key="g_email")
            g_pw = st.text_input(
                "Guest 密碼" if lang == "zh" else "Password",
                type="password",
                key="g_pw",
            )
            if st.button(
                "登入 Guest" if lang == "zh" else "Login as Guest",
                key="guest_login_btn",
            ):
                guests = load_guest_accounts()
                acc = guests.get(g_email)
                if acc and acc.get("password") == g_pw:
                    st.session_state.user_email = g_email
                    st.session_state.user_role = "free"
                    st.session_state.is_authenticated = True
                    save_state_to_disk()
                    st.experimental_rerun()
                else:
                    st.error(
                        "帳號或密碼錯誤"
                        if lang == "zh"
                        else "Invalid guest credentials"
                    )

            st.markdown("**Guest 註冊**" if lang == "zh" else "**Guest Signup**")
            new_guest_email = st.text_input(
                "註冊 Email" if lang == "zh" else "Email for signup",
                key="new_guest_email",
            )
            if st.button(
                "取得 Guest 密碼" if lang == "zh" else "Generate Guest password",
                key="guest_signup_btn",
            ):
                if not new_guest_email:
                    st.error(
                        "請輸入 Email" if lang == "zh" else "Please enter an email"
                    )
                else:
                    guests = load_guest_accounts()
                    if new_guest_email in guests or new_guest_email in ACCOUNTS:
                        st.error(
                            "Email 已存在" if lang == "zh" else "Email already exists"
                        )
                    else:
                        pw = "".join(
                            __import__("random").choice("0123456789") for _ in range(8)
                        )
                        guests[new_guest_email] = {"password": pw, "role": "free"}
                        save_guest_accounts(guests)
                        st.success(
                            f"Guest 帳號已建立！密碼：{pw}"
                            if lang == "zh"
                            else f"Guest account created! Password: {pw}"
                        )
        return

    # =========================
    # 已登入後的主頁面
    # =========================
    lang = st.session_state.lang
    user_email = st.session_state.user_email
    user_role = st.session_state.user_role
    is_guest = user_role == "free"
    model_name = resolve_model_for_user(user_role)

    title = (
        "Error-Free® 多框架 AI 文件分析"
        if lang == "zh"
        else "Error-Free® Multi-framework AI Document Analyzer"
    )
    st.title(title)
    st.markdown("---")

    doc_tracking = load_doc_tracking()

    # Step 1: Upload
    st.subheader("步驟一：上傳文件" if lang == "zh" else "Step 1: Upload document")
    uploaded = st.file_uploader(
        "請上傳 PDF / DOCX / TXT / 圖片 (JPG/PNG)"
        if lang == "zh"
        else "Upload PDF / DOCX / TXT / Image (JPG/PNG)",
        type=["pdf", "docx", "txt", "jpg", "jpeg", "png"],
    )

    if uploaded is not None:
        doc_text = read_file_to_text(uploaded)
        if doc_text:
            if is_guest:
                docs = doc_tracking.get(user_email, [])
                if len(docs) >= 3 and st.session_state.current_doc_id not in docs:
                    st.error(
                        "試用帳號最多上傳 3 份文件"
                        if lang == "zh"
                        else "Trial accounts may upload up to 3 documents only."
                    )
                else:
                    if st.session_state.current_doc_id not in docs:
                        new_id = f"doc_{datetime.datetime.now().timestamp()}"
                        docs.append(new_id)
                        doc_tracking[user_email] = docs
                        st.session_state.current_doc_id = new_id
                        save_doc_tracking(doc_tracking)
                    st.session_state.last_doc_text = doc_text
                    save_state_to_disk()
            else:
                st.session_state.current_doc_id = (
                    f"doc_{datetime.datetime.now().timestamp()}"
                )
                st.session_state.last_doc_text = doc_text
                save_state_to_disk()

    # Step 2: Framework
    st.subheader("步驟二：選擇分析框架" if lang == "zh" else "Step 2: Select framework")
    fw_keys = list(FRAMEWORKS.keys())
    fw_labels = [
        FRAMEWORKS[k]["name_zh"] if lang == "zh" else FRAMEWORKS[k]["name_en"]
        for k in fw_keys
    ]
    k2l = dict(zip(fw_keys, fw_labels))
    l2k = dict(zip(fw_labels, fw_keys))

    current_fw = st.session_state.selected_framework_key or fw_keys[0]
    selected_label = k2l[current_fw]

    new_label = st.selectbox(
        "選擇框架" if lang == "zh" else "Select framework",
        fw_labels,
        index=fw_labels.index(selected_label),
    )
    new_key = l2k[new_label]
    st.session_state.selected_framework_key = new_key

    framework_states = st.session_state.framework_states
    if new_key not in framework_states:
        framework_states[new_key] = {
            "analysis_done": False,
            "analysis_output": "",
            "followup_history": [],
            "download_used": False,
        }
    save_state_to_disk()
    current_state = framework_states[new_key]

    st.markdown("---")

    # Step 3: Run analysis
    st.subheader("步驟三：執行分析" if lang == "zh" else "Step 3: Run analysis")
    can_run = not current_state["analysis_done"]

    if can_run:
        run_btn = st.button(
            "開始分析" if lang == "zh" else "Run analysis", key="run_analysis_btn"
        )
    else:
        run_btn = False
        st.info(
            "此框架已完成一次分析"
            if lang == "zh"
            else "Analysis already completed for this framework."
        )

    if not is_guest:
        if st.button("重置（新文件）" if lang == "zh" else "Reset document"):
            st.session_state.framework_states = {}
            st.session_state.last_doc_text = ""
            st.session_state.current_doc_id = None
            save_state_to_disk()
            st.experimental_rerun()

    if run_btn and can_run:
        if not st.session_state.last_doc_text:
            st.error(
                "請先上傳文件" if lang == "zh" else "Please upload a document first."
            )
        else:
            with st.spinner("分析中..." if lang == "zh" else "Running analysis..."):
                analysis_text = run_llm_analysis(
                    new_key,
                    lang,
                    st.session_state.last_doc_text,
                    model_name,
                )
            current_state["analysis_done"] = True
            current_state["analysis_output"] = clean_report_text(analysis_text)
            current_state["followup_history"] = []
            save_state_to_disk()
            record_usage(user_email, new_key, "analysis")
            st.success("分析完成！" if lang == "zh" else "Analysis completed!")

    # Step 4: show all framework results
    any_analysis = any(s.get("analysis_output") for s in framework_states.values())

    for fw_key in FRAMEWORKS.keys():
        state = framework_states.get(fw_key)
        if not state or not state.get("analysis_output"):
            continue

        st.markdown("---")
        fw = FRAMEWORKS[fw_key]
        title_fw = (
            f"{fw['name_zh']}：分析與問答"
            if lang == "zh"
            else f"{fw['name_en']}: Analysis & Q&A"
        )
        st.subheader(("⭐ " if fw_key == new_key else "") + title_fw)

        st.markdown("#### 分析結果" if lang == "zh" else "#### Analysis result")
        st.markdown(state["analysis_output"])

        st.markdown(
            "#### 後續提問（Q&A）" if lang == "zh" else "#### Follow-up Q&A history"
        )
        if state["followup_history"]:
            for i, (q, a) in enumerate(state["followup_history"], start=1):
                st.markdown(f"**Q{i}:** {q}")
                st.markdown(f"**A{i}:** {a}")
                st.markmarkdown("---")
        else:
            st.info("尚無追問" if lang == "zh" else "No follow-up questions yet.")

        # Download
        st.markdown("##### 下載報告" if lang == "zh" else "##### Download report")
        st.caption(
            "報告只包含分析與 Q&A，不含原始文件。"
            if lang == "zh"
            else "Report includes analysis + Q&A only (no original document)."
        )

        if is_guest and state.get("download_used"):
            st.error(
                "已達下載次數上限（1 次）"
                if lang == "zh"
                else "Download limit reached (1 time)."
            )
        else:
            report = build_full_report(lang, fw_key, state)
            now_str = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

            with st.expander("Download"):
                fmt = st.radio(
                    "選擇格式" if lang == "zh" else "Select format",
                    ["Word (DOCX)", "PDF", "PowerPoint (PPTX)"],
                    key=f"fmt_{fw_key}",
                )

                data: bytes
                mime: str
                ext: str

                if fmt.startswith("Word"):
                    data = build_docx_bytes(report)
                    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    ext = "docx"
                elif fmt.startswith("PDF"):
                    data = build_pdf_bytes(report)
                    mime = "application/pdf"
                    ext = "pdf"
                else:
                    if not HAS_PPTX:
                        st.error(
                            "伺服器尚未安裝 python-pptx，無法匯出 PPTX。"
                            if lang == "zh"
                            else "python-pptx is not installed; PPTX export disabled."
                        )
                        data = b""
                        mime = "application/octet-stream"
                        ext = "pptx"
                    else:
                        try:
                            data = build_pptx_bytes(report)
                            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            ext = "pptx"
                        except Exception as e:
                            st.error(f"PPTX 匯出失敗：{e}")
                            data = b""
                            mime = "application/octet-stream"
                            ext = "pptx"

                if data:
                    clicked = st.download_button(
                        "開始下載" if lang == "zh" else "Download",
                        data=data,
                        file_name=f"errorfree_{fw_key}_{now_str}.{ext}",
                        mime=mime,
                        key=f"dl_{fw_key}_{ext}",
                    )
                    if clicked:
                        state["download_used"] = True
                        save_state_to_disk()
                        record_usage(user_email, fw_key, "download")

    # Global follow-up area
    if any_analysis:
        st.markdown("---")
        st.subheader("後續提問" if lang == "zh" else "Follow-up questions")

        curr_state = framework_states[new_key]
        if is_guest and len(curr_state["followup_history"]) >= 3:
            st.error(
                "已達追問上限（3 次）"
                if lang == "zh"
                else "Follow-up limit reached (3 times)."
            )
        else:
            extra_file = st.file_uploader(
                "上傳附加文件（可選）"
                if lang == "zh"
                else "Upload supplementary file (optional)",
                type=["pdf", "docx", "txt", "jpg", "jpeg", "png"],
                key=f"extra_{new_key}",
            )
            extra_text = read_file_to_text(extra_file) if extra_file else ""

            prompt = st.chat_input(
                f"針對 {FRAMEWORKS[new_key]['name_zh']} 的追問"
                if lang == "zh"
                else f"Ask a follow-up about {FRAMEWORKS[new_key]['name_en']}"
            )
            if prompt:
                with st.spinner("思考中..." if lang == "zh" else "Thinking..."):
                    answer = run_followup_qa(
                        new_key,
                        lang,
                        st.session_state.last_doc_text,
                        curr_state["analysis_output"],
                        prompt,
                        model_name,
                        extra_text,
                    )
                curr_state["followup_history"].append(
                    (prompt, clean_report_text(answer))
                )
                save_state_to_disk()
                record_usage(user_email, new_key, "followup")
                st.experimental_rerun()

    save_state_to_disk()


if __name__ == "__main__":
    main()
